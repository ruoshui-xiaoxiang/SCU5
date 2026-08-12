# -*- coding: utf-8 -*-
"""SCU3.0 架构与扩展方案 PPT 生成器"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 颜色方案（与 UI 蓝本一致：黑白主色 + 靛蓝强调）
COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x11, 0x11, 0x11)
COLOR_MUTED = RGBColor(0x6B, 0x72, 0x80)
COLOR_PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
COLOR_PRIMARY_SOFT = RGBColor(0xEE, 0xF2, 0xFF)
COLOR_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
COLOR_ACCENT = RGBColor(0x63, 0x66, 0xF1)
COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)
COLOR_DANGER = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height


def add_slide(title, layout_idx=6):
    """添加空白幻灯片 + 标题"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # 标题栏
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT
    return slide


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 bold=False, color=COLOR_TEXT, align=PP_ALIGN.LEFT, line_space=1.5):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(font_size * (line_space - 1))
    return txBox


def add_card(slide, left, top, width, height, title, content, accent=COLOR_PRIMARY):
    """添加卡片（带左边框强调色）"""
    # 卡片背景
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    card.line.color.rgb = COLOR_BORDER
    card.line.width = Pt(0.75)
    # 左边框强调
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()
    # 标题
    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.3), Inches(0.4),
                 title, font_size=14, bold=True, color=accent)
    # 内容
    add_text_box(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.3), height - Inches(0.6),
                 content, font_size=11, color=COLOR_TEXT, line_space=1.3)


def add_table(slide, left, top, width, height, data, header_color=COLOR_PRIMARY):
    """添加表格"""
    rows = len(data)
    cols = len(data[0]) if data else 0
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.alignment = PP_ALIGN.LEFT
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                else:
                    p.font.color.rgb = COLOR_TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA) if r % 2 == 1 else RGBColor(0xFF, 0xFF, 0xFF)
    return table_shape


# ============ 幻灯片 1：封面 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
# 背景
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid()
bg.fill.fore_color.rgb = COLOR_BG
bg.line.fill.background()
# 顶部装饰条
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.3))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = COLOR_PRIMARY
top_bar.line.fill.background()
# 底部装饰条
bot_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SH - Inches(0.3), SW, Inches(0.3))
bot_bar.fill.solid()
bot_bar.fill.fore_color.rgb = COLOR_PRIMARY
bot_bar.line.fill.background()
# 主标题
add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
             "SCU3.0 架构与扩展方案", font_size=44, bold=True, color=COLOR_TEXT, align=PP_ALIGN.CENTER)
# 副标题
add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.6),
             "三维度分离架构 · 阴阳双签 · 分布式扩展", font_size=20, color=COLOR_MUTED, align=PP_ALIGN.CENTER)
# 版本信息
add_text_box(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
             "版本：SCU3.0  ·  架构：v3  ·  日期：2026-08-11",
             font_size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# ============ 幻灯片 2：整体架构概览 ============
slide = add_slide("一、整体架构概览")
# 三维度说明
add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
             "CUF（Compute Unit Fabric）三维度分离架构：将三个正交维度彻底解耦",
             font_size=16, bold=True, color=COLOR_TEXT)
# 三张卡片
add_card(slide, Inches(0.6), Inches(1.8), Inches(3.9), Inches(1.8),
         "数据流", "感知(W2) → 记忆(W1) → 执行(W1)\n→ 认知(M) → 元认知(M) → 输出", COLOR_PRIMARY)
add_card(slide, Inches(4.7), Inches(1.8), Inches(3.9), Inches(1.8),
         "依赖方向", "D ← M ← W1 ← W2\nA4 公理：依赖方向不可反向\n（数据流不受 A4 约束）", COLOR_ACCENT)
add_card(slide, Inches(8.8), Inches(1.8), Inches(3.9), Inches(1.8),
         "守卫横切", "5 个守卫点横切在数据流管道上\n同层免审，跨层必审", COLOR_SUCCESS)
# 分层结构图
add_text_box(slide, Inches(0.6), Inches(3.9), Inches(12), Inches(0.4),
             "分层结构", font_size=16, bold=True, color=COLOR_TEXT)
layers = [
    ("W2 层", "感知入口 · 意图识别（12+ 种）· 领域识别", COLOR_ACCENT),
    ("W1 层", "工作层（记忆+执行，同层免审）· 三级记忆 · 14 种工具 · 熵税账本 · RAG 知识库", COLOR_PRIMARY),
    ("M 层", "元认知/认知层 · 多策略 LLM 推理 · 阴阳对子 · 插件市场 · 分布式执行", COLOR_SUCCESS),
    ("D 层", "基线层（只读）· 四公理 · 四契约 · 账本抽象基类 · MANIFEST 哈希基线", COLOR_DANGER),
]
for i, (name, desc, color) in enumerate(layers):
    top = Inches(4.4 + i * 0.65)
    # 层标签
    label = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(1.2), Inches(0.55))
    label.fill.solid()
    label.fill.fore_color.rgb = color
    label.line.fill.background()
    ltf = label.text_frame
    ltf.margin_top = Inches(0.08)
    lp = ltf.paragraphs[0]
    lp.text = name
    lp.font.size = Pt(13)
    lp.font.bold = True
    lp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    lp.alignment = PP_ALIGN.CENTER
    # 层描述
    add_text_box(slide, Inches(1.9), top, Inches(10.8), Inches(0.55),
                 desc, font_size=12, color=COLOR_TEXT)

# ============ 幻灯片 3：核心设计原则 ============
slide = add_slide("核心设计原则")
principles = [
    ("D 层只读", "代码定义在 D 层（只读），运行时状态在 W1 层\n消除\"审计 D 层账本需写 D 层\"的自指死循环", COLOR_PRIMARY),
    ("同层免审", "W1→W1、M→M 同层流动免审\n只有跨 CUF 层才审计，守卫点从 4 减至 3+2", COLOR_ACCENT),
    ("阴阳双签", "阴方 DeepSeek（批判）+ 阳方 Qwen（支持）+ 太极合一\nγ_yin≥0.75 / γ_yang≥0.65", COLOR_SUCCESS),
    ("分布式+插件闭环", "能力缺失 → 自动下载加载 → 重试 → 经验沉淀 → 自进化\n失败经验 ≥3 次触发异步自进化扫描", COLOR_DANGER),
    ("多重安全", "API Key 时序防护 · SANDBOX 沙箱 · AST 预检\nD 层完整性熔断 · 50+ 内容过滤规则", COLOR_PRIMARY),
    ("扩展能力已就绪", "横向加 Worker · 纵向多单元\n阴阳对子可叠加在 Worker 内", COLOR_ACCENT),
]
for i, (title, desc, color) in enumerate(principles):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.3 + row * 2.8)
    add_card(slide, left, top, Inches(3.9), Inches(2.5), title, desc, color)

# ============ 幻灯片 4：守护系统 ============
slide = add_slide("三、守护系统 — 5 个守卫点")
# 守卫点表格
guard_data = [
    ["#", "守卫点", "触发条件", "实现文件"],
    ["①", "W2→W1 跨层审计", "跨 CUF 层", "firewall.py"],
    ["②", "W1→M 跨层审计", "跨 CUF 层", "firewall.py"],
    ["③", "工具守卫", "无论同层与否", "tool_guard.py"],
    ["④", "周期审计", "M→W1 同层免审", "metacognition.py"],
    ["⑤", "内容过滤", "响应生成后", "content_filter.py"],
]
add_table(slide, Inches(0.6), Inches(1.2), Inches(7.5), Inches(3.5), guard_data)
# CUF 防火墙审计顺序
add_text_box(slide, Inches(8.4), Inches(1.2), Inches(4.5), Inches(0.4),
             "CUF 防火墙审计顺序", font_size=14, bold=True, color=COLOR_PRIMARY)
audit_steps = "1. 层标识符校验\n2. 同层免审短路\n3. 白名单短路\n4. A1 基线不可变性\n5. A4 层级单向性\n6. A3 契约闭环性\n7. A2 熵税经济性"
add_text_box(slide, Inches(8.4), Inches(1.7), Inches(4.5), Inches(3),
             audit_steps, font_size=12, color=COLOR_TEXT, line_space=1.6)
# 底部说明
add_card(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.8),
         "D 层完整性校验",
         "· 启动校验：失败时熔断拒绝启动\n· 定期校验：每小时一次，失败时进入只读降级模式\n· 基线 hash：从 MANIFEST.expected_hashes 加载\n· 禁止运行时状态：扫描 _balance/_history/threading.Lock() 等模式",
         COLOR_DANGER)

# ============ 幻灯片 5：阴阳对子思考 ============
slide = add_slide("四、阴阳对子思考（方案 C）")
# 触发条件
add_card(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.0),
         "触发条件",
         "intent == \"analytical\" 且无工具结果（正则覆盖：分析/批判/反思/利弊/可行性/前景/影响/优劣）",
         COLOR_ACCENT)
# 三方角色
roles = [
    ("阴方", "DeepSeek-Chat", "批判视角", "找漏洞、风险、反对理由\n至少 3 条具体反对意见", COLOR_PRIMARY),
    ("阳方", "Qwen-Plus", "支持视角", "找优势、机会、支持理由\n失败回退 DeepSeek", COLOR_ACCENT),
    ("合一", "DeepSeek-Chat", "综合视角", "不复读素材、不提及阴方阳方\n独立观点、500-800 字", COLOR_SUCCESS),
]
for i, (role, llm, view, prompt, color) in enumerate(roles):
    left = Inches(0.6 + i * 4.1)
    add_card(slide, left, Inches(2.5), Inches(3.9), Inches(2.5), f"{role} · {llm}", f"{view}\n\n{prompt}", color)
# 评分维度
add_text_box(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(0.4),
             "评分维度（满分 1.0）", font_size=14, bold=True, color=COLOR_PRIMARY)
score_data = [
    ["维度", "分值", "说明"],
    ["基础分", "0.3", "起始分"],
    ["批判/支持性词汇", "+0.2", "词汇命中"],
    ["分点论证", "+0.15/+0.05/+0.05", "三点递减"],
    ["因果论证", "+0.15", "因果链"],
    ["字数充分", "+0.2/+0.1", ">200 / >100"],
]
add_table(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.5), score_data)

# ============ 幻灯片 6：认知层处理流程 ============
slide = add_slide("四、认知层处理流程")
add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
             "AND 关系非互斥 — 多策略综合注入 LLM", font_size=16, bold=True, color=COLOR_TEXT)
# 流程卡片
flows = [
    ("analytical 意图", "阴阳对子思考\n阴DeepSeek + 阳Qwen + 合一\n成功→返回 / 失败→降级", COLOR_PRIMARY),
    ("web_search 成功", "搜索结果 + 深度爬取\n+ RAG 综合注入 LLM", COLOR_ACCENT),
    ("web_crawl 成功", "爬取结果 + RAG\n综合注入 LLM", COLOR_SUCCESS),
    ("其他工具成功", "13 种工具格式化输出\n信息查询类额外注入 RAG", COLOR_PRIMARY),
    ("工具全失败", "插件市场闭环 → 仍失败\n→ RAG + LLM 常规对话", COLOR_DANGER),
    ("web_search 意图无工具", "兜底联网搜索\n避免 LLM 凭空说不能联网", COLOR_ACCENT),
    ("无工具调用", "RAG 上下文 + LLM 生成回复\n（闲聊也注入 RAG）", COLOR_SUCCESS),
]
for i, (title, desc, color) in enumerate(flows):
    col = i % 4
    row = i // 4
    left = Inches(0.6 + col * 3.1)
    top = Inches(1.8 + row * 2.7)
    add_card(slide, left, top, Inches(2.9), Inches(2.4), title, desc, color)

# ============ 幻灯片 7：分布式执行 ============
slide = add_slide("五、分布式执行")
# 核心组件
add_text_box(slide, Inches(0.6), Inches(1.1), Inches(6), Inches(0.4),
             "核心组件", font_size=14, bold=True, color=COLOR_PRIMARY)
comp_data = [
    ["组件", "职责"],
    ["WorkerNode", "工作节点（IDLE/BUSY/OFFLINE，能力声明 cpu/memory/gpu/special_tools）"],
    ["WorkerRegistry", "节点注册表（轮询/最少忙碌/能力匹配）"],
    ["TaskDispatcher", "任务分发器（七种状态，幂等缓存，重试+超时迁移）"],
    ["WorkerServer", "工作节点服务端（5 个 HTTP 端点）"],
    ["LocalMultiProcessExecutor", "本地多进程降级（multiprocessing.Pool）"],
    ["DistributedExecutor", "主类（自动选择模式，状态持久化）"],
]
add_table(slide, Inches(0.6), Inches(1.6), Inches(7.5), Inches(4), comp_data)
# 负载均衡策略
add_text_box(slide, Inches(8.4), Inches(1.1), Inches(4.5), Inches(0.4),
             "负载均衡策略", font_size=14, bold=True, color=COLOR_PRIMARY)
strategies = "ROUND_ROBIN（轮询）\n\nLEAST_BUSY（最少失败+最少任务）\n\nCAPABILITY_MATCH（能力最贴近\n  选资源最少占用的\n  留大节点给重任务）"
add_text_box(slide, Inches(8.4), Inches(1.6), Inches(4.5), Inches(2.5),
             strategies, font_size=12, color=COLOR_TEXT, line_space=1.5)
# 故障处理
add_card(slide, Inches(8.4), Inches(4.3), Inches(4.5), Inches(2.8),
         "故障处理",
         "· 心跳超时 30s → OFFLINE\n· 任务迁移 retry()\n· 幂等性：重复 task_id 跳过\n· 状态持久化：\n  SCU3_data/distributed_state.json\n· 重启后恢复（标记 offline 待心跳确认）",
         COLOR_DANGER)
# 分片合并
add_text_box(slide, Inches(0.6), Inches(5.8), Inches(7.5), Inches(0.4),
             "任务分片：list 均分 / dict 按 key / 不可分片复制冗余  ·  合并：concat/sum/avg/max/min/dict_merge/first",
             font_size=11, color=COLOR_MUTED)

# ============ 幻灯片 8：插件系统 ============
slide = add_slide("六、插件系统")
# 能力匹配
add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
             "能力匹配（四级优先级）", font_size=14, bold=True, color=COLOR_PRIMARY)
match_data = [
    ["优先级", "匹配方式", "示例"],
    ["1", "文件扩展名匹配", ".pdf→pdf_reader / .docx→docx_reader / .xlsx→excel_reader"],
    ["2", "触发词匹配", "市场清单 triggers 字段"],
    ["3", "能力关键词匹配", "市场清单 capabilities 字段"],
    ["4", "失败工具→能力映射", "tool_capability_map 转插件名"],
]
add_table(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(2.2), match_data)
# 生命周期
add_text_box(slide, Inches(0.6), Inches(4.1), Inches(6), Inches(0.4),
             "生命周期", font_size=14, bold=True, color=COLOR_PRIMARY)
add_text_box(slide, Inches(0.6), Inches(4.6), Inches(6), Inches(2.5),
             "· 安装：pip / git 两种方式\n  pip 多源回退：清华→阿里云→官方\n  git：GitHub 白名单 + clone --depth 1\n· 加载：importlib 动态导入 + 工具工厂\n· TTL：默认 600s，每 30s 检查\n· 卸载：注销工具 + 卸载 Python 模块\n· 持久化：keep_alive() 标记持久模式",
             font_size=11, color=COLOR_TEXT, line_space=1.4)
# 经验沉淀
add_card(slide, Inches(6.8), Inches(4.1), Inches(5.9), Inches(3),
         "经验沉淀与自进化",
         "· 成功路径记录，下次直接预加载跳过 all_failed\n· 衰减机制：30 天未用降权\n· 成熟阈值：成功 2 次以上视为成熟方案\n· 自进化触发：fail_count ≥ 3 且 success_count == 0\n· 异步触发，不阻塞用户响应\n· 闭环：缺陷分析→提案→双签→审批→应用+备份",
         COLOR_ACCENT)

# ============ 幻灯片 9：扩展方案 - 横向扩展 ============
slide = add_slide("八、扩展方案 — 横向扩展（添加 Worker 节点）")
# 方式 A
add_card(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(2.8),
         "方式 A：添加远程 Worker（独立机器部署）",
         "1. 远程机器启动 WorkerServer：\n"
         "   from m_layer.distributed_executor import WorkerServer\n"
         "   server = WorkerServer(port=9700, handler=my_task_handler, host=\"0.0.0.0\")\n"
         "   server.start(background=False)\n\n"
         "2. 主节点注册远程 Worker：\n"
         "   executor = get_distributed_executor()\n"
         "   executor.add_remote_worker(url=\"http://remote-host:9700\", capabilities={\"cpu\": 8, \"memory\": 16384, \"gpu\": 1})\n"
         "   或通过 API：POST /distributed/workers/add（需 admin key）\n\n"
         "3. 任务分发：POST /distributed/execute，自动选择分布式/本地模式",
         COLOR_PRIMARY)
# 方式 B
add_card(slide, Inches(0.6), Inches(4.2), Inches(5.9), Inches(2.8),
         "方式 B：本地多进程模拟",
         "无需远程节点\nLocalMultiProcessExecutor 使用\nmultiprocessing.Pool\n默认 worker 数 = cpu_count - 1\n任务分片 split_task()\n结果合并 merge_results()",
         COLOR_ACCENT)
# 关键扩展点
add_card(slide, Inches(6.8), Inches(4.2), Inches(5.9), Inches(2.8),
         "关键扩展点",
         "· 负载均衡：round_robin / least_busy / capability_match\n· 能力声明：cpu/memory/gpu/special_tools\n· 故障处理：心跳超时 30s → OFFLINE → 任务迁移\n· 幂等性：重复 task_id 跳过分发\n· 状态持久化：distributed_state.json",
         COLOR_SUCCESS)

# ============ 幻灯片 10：扩展方案 - 多单元 + 阴阳叠加 ============
slide = add_slide("八、扩展方案 — 多单元 + 阴阳叠加")
# 多单元
add_card(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(2.8),
         "多单元：修改 /units 返回多个单元配置",
         "当前 /units 返回单个 SCU3-default。扩展方式：\n\n"
         "· 修改 /units 端点返回多个单元（不同 system_prompt_style / model / domain / force_yin_yang）\n"
         "· 前端已就绪：chatUnit 下拉框 + loadUnits() 已支持多选项\n"
         "· 后端联动：/chat 请求体增加 uid 字段，按 uid 选择配置\n"
         "· 隔离级别：\n"
         "  软隔离（当前）：共享 ledger/记忆/知识库\n"
         "  硬隔离：每个单元独立 DATA_DIR + 独立 ledger 实例",
         COLOR_PRIMARY)
# 阴阳叠加
add_card(slide, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.8),
         "阴阳对子叠加：每个 Worker 内部都可跑阴阳对子",
         "阴阳对子思考与分布式执行是正交能力，可叠加：\n\n"
         "· Worker 内嵌阴阳对子：handler 回调内调用 _yin_yang_think()，每个 Worker 可独立配置 LLM 平台\n"
         "· 分布式阴阳对子：split_task() 拆为 3 子任务（阴方/阳方/合一），分发到不同 Worker 并行执行\n"
         "· 双签判定位置：主节点统一双签（推荐）或 Worker 端本地双签\n"
         "· 安全约束不变：软双签，不触发 Pair 硬约束，不跨层，不修改 D 层",
         COLOR_SUCCESS)

# ============ 幻灯片 11：安全约束 ============
slide = add_slide("九、安全约束")
sec_items = [
    ("API Key 认证", "双 Key 体系（普通+管理员）\nsecrets.compare_digest() 防时序侧信道\n敏感端点集合需 admin 权限", COLOR_PRIMARY),
    ("文件操作限制", "SANDBOX_DIR 隔离\n_safe_path() 用 commonpath 防前缀碰撞\n沙箱执行：AST 预检 + 5s 超时", COLOR_ACCENT),
    ("代码自修改保护", "D 层保护清单 + 危险模式黑名单\n文件扩展名白名单 {.py}\n阴阳双签 + 人工审批 + 备份回滚", COLOR_DANGER),
    ("内容过滤", "50+ 条正则规则\nAPI Key/密码/内网IP/数据库连接串\nJWT/手机号/身份证/邮箱", COLOR_SUCCESS),
    ("模块保护", "PROTECTED_MODULES 不可卸载\nfirewall/entropy_ledger/axioms\nengine/meta_guard/code_self_modify", COLOR_PRIMARY),
    ("网络与监听", "默认监听 127.0.0.1\n0.0.0.0 告警\n默认端口 8300", COLOR_ACCENT),
]
for i, (title, desc, color) in enumerate(sec_items):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.3 + row * 2.8)
    add_card(slide, left, top, Inches(3.9), Inches(2.5), title, desc, color)

# ============ 幻灯片 12：API 接口分类 ============
slide = add_slide("十、API 接口分类清单（100+ 端点）")
api_data = [
    ["分类", "数量", "主要端点"],
    ["对话与会话", "8", "/chat · /chat/stream · /chat/image · /conversation/*"],
    ["CUF 守卫与状态", "10", "/health · /status · /pair/status · /cognition/yin-yang · /cuf/*"],
    ["分布式执行", "7", "/distributed/execute · /distributed/workers/* · /distributed/health"],
    ["自修改与自进化", "15", "/self-modify/* · /evolution/* · /learning/* · /code/proposals"],
    ["插件与市场", "15", "/plugins/* · /plugins/market/*（install/unload/uninstall/match）"],
    ["知识库与向量", "8", "/knowledge/* · /vector/*（search/migrate）"],
    ["三级记忆", "7", "/memory/*（stats/health/search/episode/knowledge）"],
    ["LLM 与本地模型", "10+", "/llm/* · /models · /units · /local-model/* · /vision/*"],
    ["自动化与浏览器", "15+", "/automation/* · /browser/*（start/navigate/click/fill/screenshot）"],
    ["多模态与语音", "7", "/multimodal/* · /voice/*"],
    ["MCP 协议", "6", "/mcp/*（tools/call/connect/servers/health）"],
    ["模块注册表", "5", "/modules/*（load/unload/enable/disable/status）"],
]
add_table(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.8), api_data)

# ============ 幻灯片 13：总结 ============
slide = add_slide("十一、总结")
# 架构亮点
add_text_box(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
             "当前架构亮点", font_size=16, bold=True, color=COLOR_PRIMARY)
highlights = [
    ("D 层只读", "代码定义与运行时状态分离，消除自指死循环"),
    ("同层免审", "W1→W1、M→M 免审，跨层才审计"),
    ("阴阳双签", "阴DeepSeek + 阳Qwen + 合一，γ_yin≥0.75/γ_yang≥0.65"),
    ("分布式+插件闭环", "能力缺失→自动下载→重试→经验沉淀→自进化"),
    ("多重安全", "API Key 时序防护 · 沙箱 · AST 预检 · 熔断 · 内容过滤"),
]
for i, (title, desc) in enumerate(highlights):
    top = Inches(1.7 + i * 0.55)
    add_text_box(slide, Inches(0.8), top, Inches(3), Inches(0.5),
                 f"✓ {title}", font_size=13, bold=True, color=COLOR_SUCCESS)
    add_text_box(slide, Inches(3.8), top, Inches(9), Inches(0.5),
                 desc, font_size=12, color=COLOR_TEXT)
# 扩展能力
add_text_box(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.4),
             "扩展能力已就绪", font_size=16, bold=True, color=COLOR_ACCENT)
extends = [
    ("横向扩展", "添加 Worker 节点（远程/本地多进程）", COLOR_PRIMARY),
    ("纵向扩展", "多单元配置（修改 /units 返回值）", COLOR_ACCENT),
    ("能力叠加", "阴阳对子可叠加在 Worker 内（分布式+双签）", COLOR_SUCCESS),
]
for i, (title, desc, color) in enumerate(extends):
    left = Inches(0.6 + i * 4.1)
    add_card(slide, left, Inches(5.2), Inches(3.9), Inches(1.8), title, desc, color)
# 结束语
add_text_box(slide, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4),
             "所有扩展点都有对应的 API 端点和单例管理器，架构弹性充足。",
             font_size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# ============ 保存 ============
desktop = os.path.join(os.environ['USERPROFILE'], 'Desktop')
output_path = os.path.join(desktop, 'SCU3.0架构与扩展方案.pptx')
prs.save(output_path)
print(f"PPT 已保存：{output_path}")
print(f"共 {len(prs.slides)} 张幻灯片")
